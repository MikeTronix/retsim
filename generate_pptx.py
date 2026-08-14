import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide width and height to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6] # Blank slide

    # Color Palette
    DARK_NAVY = RGBColor(15, 23, 42)      # #0F172A
    SLATE_BLUE = RGBColor(30, 41, 59)     # #1E293B
    EMERALD = RGBColor(16, 185, 129)      # #10B981
    LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF
    TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A
    TEXT_MUTED = RGBColor(100, 116, 139)  # #64748B
    BORDER_COLOR = RGBColor(226, 232, 240) # #E2E8F0

    def add_header(slide, title_text, category_text="RETSIM ARCHITECTURE"):
        # Top banner category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = EMERALD

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(24)
        p_title.font.bold = True
        p_title.font.color.rgb = TEXT_DARK

    def set_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # ==========================================
    # SLIDE 1: Title Slide (Dark Elegant)
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    set_background(slide1, DARK_NAVY)

    # Accent decorative box
    accent = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = EMERALD
    accent.line.fill.background()

    # Title box
    tbox = slide1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(3.5))
    tf = tbox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "retsim"
    p0.font.size = Pt(44)
    p0.font.bold = True
    p0.font.color.rgb = EMERALD

    p1 = tf.add_paragraph()
    p1.text = "Institutional-Grade Retirement Simulation Engine"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(255, 255, 255)
    p1.space_before = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "Next-Generation Wealth, Tax, & Cash Flow Modeling Built for Fiduciaries & Financial Planners"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.space_before = Pt(14)

    # Footer pills on Slide 1
    foot_box = slide1.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(11.0), Inches(0.8))
    tf_foot = foot_box.text_frame
    p_foot = tf_foot.paragraphs[0]
    p_foot.text = "Monthly Calculation Granularity  •  Tiered Budget & QoL Metric  •  Medicare IRMAA & SS Haircut Engine  •  Cholesky Stochastic Returns"
    p_foot.font.size = Pt(12)
    p_foot.font.color.rgb = RGBColor(203, 213, 225)

    # ==========================================
    # SLIDE 2: Why Legacy Calculators Fail Planners
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    set_background(slide2, LIGHT_BG)
    add_header(slide2, "Beyond Static 4% Rules & Oversimplified Spreadsheets", "THE ADVISORY CHALLENGE")

    # 4 Cards Layout (2x2 Grid)
    card_data = [
        ("Path-Dependent Tax Friction", "Retail tools treat tax rates as static averages. retsim computes exact progressive federal/state brackets, Social Security combined income taxability, and 2-year lookback IRMAA surcharges."),
        ("Rigid All-or-Nothing Budgets", "Clients don't spend identically every month. retsim separates non-discretionary BasicBudget from discretionary QoLBudget, tracking cumulative QoL satisfaction (0-100%)."),
        ("Coarse Annual Time Steps", "Annual calculations miss mid-year Social Security claims, monthly mortgage principal decay, and intra-year cash flow shortfalls. retsim steps internally every month (dt = 1/12)."),
        ("Uncorrelated Risk Assumptions", "Standard tools assume stock and bond returns are independent. retsim utilizes Cholesky decomposition with PSD matrix repair to model correlated Sequence of Returns Risk (SRR)."),
    ]

    coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.5)),
        (Inches(6.8), Inches(4.5)),
    ]

    for (head, body), (x, y) in zip(card_data, coords):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.3)
        tf_c.margin_right = Inches(0.3)
        tf_c.margin_top = Inches(0.25)

        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = DARK_NAVY

        p_b = tf_c.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(8)

    # ==========================================
    # SLIDE 3: Tiered Budget & QoL Satisfaction Metric
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    set_background(slide3, LIGHT_BG)
    add_header(slide3, "Client-Centric Budget Hierarchy: BasicBudget vs. QoLBudget", "BUDGET & LIFESTYLE MODELING")

    # Left Box: Basic Budget
    box_left = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    box_left.fill.solid()
    box_left.fill.fore_color.rgb = CARD_BG
    box_left.line.color.rgb = BORDER_COLOR

    tf_l = box_left.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = Inches(0.4)
    tf_l.margin_top = Inches(0.4)

    p_l1 = tf_l.paragraphs[0]
    p_l1.text = "1. BasicBudget (Mandatory Baseline)"
    p_l1.font.size = Pt(18)
    p_l1.font.bold = True
    p_l1.font.color.rgb = DARK_NAVY

    p_l2 = tf_l.add_paragraph()
    p_l2.text = "• Core Essential Expenses: Housing, debt/mortgage, utilities, food, Medicare Part B/D premiums.\n• Priority Rule: Must always be 100% satisfied each month.\n• Solvency Guardrail: Triggers an explicit Insolvency Flag if portfolio liquidity is exhausted."
    p_l2.font.size = Pt(13)
    p_l2.font.color.rgb = TEXT_MUTED
    p_l2.space_before = Pt(14)

    # Right Box: QoL Budget
    box_right = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    box_right.fill.solid()
    box_right.fill.fore_color.rgb = CARD_BG
    box_right.line.color.rgb = BORDER_COLOR

    tf_r = box_right.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = Inches(0.4)
    tf_r.margin_top = Inches(0.4)

    p_r1 = tf_r.paragraphs[0]
    p_r1.text = "2. QoLBudget & Satisfaction Index"
    p_r1.font.size = Pt(18)
    p_r1.font.bold = True
    p_r1.font.color.rgb = EMERALD

    p_r2 = tf_r.add_paragraph()
    p_r2.text = "• Discretionary Lifestyle: Travel, dining out, hobbies, family gifting.\n• Delivery Rule: Satisfied to maximum extent possible after BasicBudget & taxes are met.\n• Cumulative Satisfaction Index (0–100%):"
    p_r2.font.size = Pt(13)
    p_r2.font.color.rgb = TEXT_MUTED
    p_r2.space_before = Pt(14)

    # Math Callout Box inside Right Box
    math_box = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.3), Inches(4.5), Inches(4.8), Inches(1.2))
    math_box.fill.solid()
    math_box.fill.fore_color.rgb = DARK_NAVY
    math_box.line.fill.background()

    tf_m = math_box.text_frame
    tf_m.word_wrap = True
    p_m = tf_m.paragraphs[0]
    p_m.text = "QoL Cumulative Satisfaction % =\n( Total Delivered QoL / Total Target QoL ) × 100%"
    p_m.font.size = Pt(13)
    p_m.font.bold = True
    p_m.font.color.rgb = RGBColor(255, 255, 255)
    p_m.alignment = PP_ALIGN.CENTER

    # ==========================================
    # SLIDE 4: Cash Flow Waterfall & Modular Strategy Pattern
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    set_background(slide4, LIGHT_BG)
    add_header(slide4, "Priority Waterfall & Modular Withdrawal Behaviors", "CASH FLOW ENGINE")

    # Horizontal Flow Visual for Priority Waterfall
    wf_title = slide4.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(0.4))
    tf_wft = wf_title.text_frame
    p_wft = tf_wft.paragraphs[0]
    p_wft.text = "Default Portfolio Withdrawal Priority Waterfall:"
    p_wft.font.size = Pt(14)
    p_wft.font.bold = True
    p_wft.font.color.rgb = DARK_NAVY

    wf_steps = [
        ("Step 1", "Taxable Brokerage", "Already-taxed basis\nLTCG rates"),
        ("Step 2", "401(k) Accounts", "Pre-tax growth\nOrdinary income"),
        ("Step 3", "Traditional IRA", "Pre-tax growth\nSubject to RMDs"),
        ("Step 4", "Roth IRA", "Tax-free growth\nNo RMD friction"),
    ]

    for i, (step, name, desc) in enumerate(wf_steps):
        x = Inches(0.8 + i * 3.0)
        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.1), Inches(2.7), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_NAVY if i == 0 else CARD_BG
        card.line.color.rgb = EMERALD if i == 0 else BORDER_COLOR

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.15)
        tf_c.margin_top = Inches(0.15)

        p_st = tf_c.paragraphs[0]
        p_st.text = step.upper()
        p_st.font.size = Pt(10)
        p_st.font.bold = True
        p_st.font.color.rgb = EMERALD if i != 0 else RGBColor(52, 211, 153)

        p_nm = tf_c.add_paragraph()
        p_nm.text = name
        p_nm.font.size = Pt(14)
        p_nm.font.bold = True
        p_nm.font.color.rgb = RGBColor(255, 255, 255) if i == 0 else DARK_NAVY

        p_ds = tf_c.add_paragraph()
        p_ds.text = desc
        p_ds.font.size = Pt(10)
        p_ds.font.color.rgb = RGBColor(203, 213, 225) if i == 0 else TEXT_MUTED

    # Bottom Strategy Pattern Extensions Box
    strat_box = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(4.1), Inches(11.7), Inches(2.6))
    strat_box.fill.solid()
    strat_box.fill.fore_color.rgb = CARD_BG
    strat_box.line.color.rgb = BORDER_COLOR

    tf_sb = strat_box.text_frame
    tf_sb.word_wrap = True
    tf_sb.margin_left = Inches(0.3)
    tf_sb.margin_top = Inches(0.25)

    p_sbt = tf_sb.paragraphs[0]
    p_sbt.text = "Pluggable WithdrawalBehavior Protocol Strategy Hooks:"
    p_sbt.font.size = Pt(16)
    p_sbt.font.bold = True
    p_sbt.font.color.rgb = DARK_NAVY

    strategies = [
        ("Dynamic Spending Guardrails (Guyton-Klinger / VPW):", "Automatically scales QoLBudget down during market drawdowns to protect portfolio longevity, and raises spending during prosperity years."),
        ("Systematic Pre-RMD Roth Conversions:", "Executes tactical Traditional IRA to Roth conversions during low-income retirement gap years to fill low tax brackets and shrink future RMDs."),
        ("Tax-Aware Portfolio Rebalancing:", "Directs monthly cash withdrawals from over-weighted asset classes to restore target asset allocation (e.g. 60/40) without triggering unnecessary tax friction."),
    ]

    for title, detail in strategies:
        p_st = tf_sb.add_paragraph()
        p_st.text = f"• {title} {detail}"
        p_st.font.size = Pt(12)
        p_st.font.color.rgb = TEXT_MUTED
        p_st.space_before = Pt(6)

    # ==========================================
    # SLIDE 5: Real-World Friction: Taxes, IRMAA, SS
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    set_background(slide5, LIGHT_BG)
    add_header(slide5, "Real-World Friction: Tax Brackets, IRMAA Cliffs, & SS Haircuts", "TAX & REGULATORY PRECISION")

    friction_cards = [
        ("Medicare Part B/D & IRMAA Lookback", "Evaluates MAGI from 2 years prior (T-2) to compute tiered IRMAA surcharges. Prevents clients from accidentally tripping steep Medicare cliff penalties via large RMDs or lump-sum withdrawals."),
        ("Social Security 'Combined Income' Tax", "Models unindexed Federal Provisional Income thresholds ($25k/$32k and $34k/$44k). Captures the expanding tax drag on Social Security benefits over multi-decade retirement timelines."),
        ("SS Insolvency Haircut Cliff", "Programmable Trustee Report benefit cuts (e.g. 23% reduction starting in 2033) while preserving COLA compounding. Allows advisors to stress-test Social Security Trust Fund depletion scenarios."),
        ("Modular Alabama State Tax Engine", "Includes AL state tax rules: 100% exemption for Social Security and Defined Benefit Pensions, 2%/4%/5% rate brackets, and AL standard deduction phaseouts."),
    ]

    coords = [
        (Inches(0.8), Inches(1.8)),
        (Inches(6.8), Inches(1.8)),
        (Inches(0.8), Inches(4.5)),
        (Inches(6.8), Inches(4.5)),
    ]

    for (head, body), (x, y) in zip(friction_cards, coords):
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.7), Inches(2.3))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.3)
        tf_c.margin_top = Inches(0.25)

        p_h = tf_c.paragraphs[0]
        p_h.text = head
        p_h.font.size = Pt(16)
        p_h.font.bold = True
        p_h.font.color.rgb = DARK_NAVY

        p_b = tf_c.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(12)
        p_b.font.color.rgb = TEXT_MUTED
        p_b.space_before = Pt(8)

    # ==========================================
    # SLIDE 6: Correlated Stochastic Returns Engine
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    set_background(slide6, LIGHT_BG)
    add_header(slide6, "Correlated Asset Classes & Sequence of Returns Risk", "STOCHASTIC ENGINE")

    # Left Side: Explanation Box
    box_l = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(6.0), Inches(4.9))
    box_l.fill.solid()
    box_l.fill.fore_color.rgb = CARD_BG
    box_l.line.color.rgb = BORDER_COLOR

    tf_bl = box_l.text_frame
    tf_bl.word_wrap = True
    tf_bl.margin_left = Inches(0.3)
    tf_bl.margin_top = Inches(0.3)

    p_bl1 = tf_bl.paragraphs[0]
    p_bl1.text = "Multivariate Return Matrix via Cholesky Factorization"
    p_bl1.font.size = Pt(16)
    p_bl1.font.bold = True
    p_bl1.font.color.rgb = DARK_NAVY

    points = [
        ("Up to 10 Asset Categories:", "Tracks equities, fixed income, cash, real estate with specific gross returns, fees, and tax categories."),
        ("Cross-Correlation Matrix (Σ = L Lᵀ):", "Simulates realistic multi-asset co-movement using Cholesky factor L applied to standard normal vectors z."),
        ("Automated PSD Matrix Repair:", "User-defined correlation matrices often fail Positive Semi-Definiteness. retsim incorporates Higham's eigenvalue clipping algorithm to guarantee mathematical validity."),
        ("3-Stage Architecture Roadmap:", "(a) Deterministic Engine → (b) Correlated Stochastic Engine → (c) Vectorized Monte Carlo Simulation."),
    ]

    for title, body in points:
        p_pt = tf_bl.add_paragraph()
        p_pt.text = f"• {title} {body}"
        p_pt.font.size = Pt(12)
        p_pt.font.color.rgb = TEXT_MUTED
        p_pt.space_before = Pt(10)

    # Right Side: Math & Matrix Card
    box_r = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(1.8), Inches(5.4), Inches(4.9))
    box_r.fill.solid()
    box_r.fill.fore_color.rgb = DARK_NAVY
    box_r.line.fill.background()

    tf_br = box_r.text_frame
    tf_br.word_wrap = True
    tf_br.margin_left = Inches(0.4)
    tf_br.margin_top = Inches(0.4)

    p_br1 = tf_br.paragraphs[0]
    p_br1.text = "Stochastic Formulation"
    p_br1.font.size = Pt(18)
    p_br1.font.bold = True
    p_br1.font.color.rgb = EMERALD

    p_br2 = tf_br.add_paragraph()
    p_br2.text = "Monthly Correlated Return Vector:\nr_m = μ_m + L · z_m\n\nWhere:\n• μ_m = Expected monthly net return vector\n• L = Lower triangular Cholesky factor of Σ\n• z_m ~ N(0, I) = Independent standard normal vector\n• Σ = User cross-correlation & covariance matrix"
    p_br2.font.size = Pt(13)
    p_br2.font.color.rgb = RGBColor(255, 255, 255)
    p_br2.space_before = Pt(14)

    p_br3 = tf_br.add_paragraph()
    p_br3.text = "Accurately captures Sequence of Returns Risk (SRR) in early retirement years."
    p_br3.font.size = Pt(12)
    p_br3.font.italic = True
    p_br3.font.color.rgb = RGBColor(52, 211, 153)
    p_br3.space_before = Pt(20)

    # ==========================================
    # SLIDE 7: Institutional Verification & Auditability
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    set_background(slide7, LIGHT_BG)
    add_header(slide7, "3-Tier Verification Suite & Institutional Auditability", "TESTING & QUALITY ASSURANCE")

    tiers = [
        ("Tier 1: Neutral Portfolio Unit Tests", "Analytical Mathematical Verification", [
            "Zero-Tax / Zero-Inflation pure compound growth verification.",
            "Mortgage amortization payoff to exactly $0.00 at month 360.",
            "Pure RMD decay matching IRS Publication 590-B factors.",
            "Boundary cliff tests at MAGI threshold ± $1."
        ]),
        ("Tier 2: External Benchmark Parity", "Gold-Standard Cross-Validation", [
            "IRS Form 1040 tax instruction table parity tests.",
            "Regression checks against baseline FireCalc / CFiresim CSV ledgers."
        ]),
        ("Tier 3: Property-Based Invariants", "Automated Hypothesis Framework", [
            "Conservation of Funds: Total Inflow = Outflow + Delta Wealth.",
            "Monotonic Tax: Higher gross income never decreases tax liability.",
            "Stochastic covariance convergence within Frobenius norm < 0.005."
        ]),
    ]

    for i, (t_title, t_sub, t_bullets) in enumerate(tiers):
        x = Inches(0.8 + i * 4.0)
        card = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(3.7), Inches(4.9))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.25)
        tf_c.margin_top = Inches(0.25)

        p_tt = tf_c.paragraphs[0]
        p_tt.text = t_title
        p_tt.font.size = Pt(15)
        p_tt.font.bold = True
        p_tt.font.color.rgb = DARK_NAVY

        p_sub = tf_c.add_paragraph()
        p_sub.text = t_sub
        p_sub.font.size = Pt(11)
        p_sub.font.bold = True
        p_sub.font.color.rgb = EMERALD
        p_sub.space_before = Pt(4)

        for bullet in t_bullets:
            p_b = tf_c.add_paragraph()
            p_b.text = f"• {bullet}"
            p_b.font.size = Pt(11)
            p_b.font.color.rgb = TEXT_MUTED
            p_b.space_before = Pt(8)

    # ==========================================
    # SLIDE 8: Summary for Financial Advisors
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    set_background(slide8, DARK_NAVY)

    # Title box
    tbox8 = slide8.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.7), Inches(1.2))
    tf8 = tbox8.text_frame
    tf8.word_wrap = True
    p8_t = tf8.paragraphs[0]
    p8_t.text = "Why retsim Transforms Retirement Advisory"
    p8_t.font.size = Pt(28)
    p8_t.font.bold = True
    p8_t.font.color.rgb = RGBColor(255, 255, 255)

    p8_sub = tf8.add_paragraph()
    p8_sub.text = "Elevating client conversations from black-box estimates to rigorous fiduciary planning."
    p8_sub.font.size = Pt(15)
    p8_sub.font.color.rgb = RGBColor(148, 163, 184)
    p8_sub.space_before = Pt(6)

    # 3 Summary Pillar Cards
    pillars = [
        ("Fiduciary Rigor", "Replaces black-box spreadsheets with an open-source, fully typed Python framework validated by 13 automated unit tests."),
        ("Actionable Client Metric", "Quantifies client lifestyle satisfaction (0-100% QoL Index) alongside portfolio survival, enabling clear trade-off discussions."),
        ("Tax Optimization Engine", "Stress-tests tactical Roth conversions, Social Security claiming ages, and 2-year IRMAA lookback friction in seconds."),
    ]

    for i, (p_title, p_body) in enumerate(pillars):
        x = Inches(0.8 + i * 4.0)
        card = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.2), Inches(3.7), Inches(4.3))
        card.fill.solid()
        card.fill.fore_color.rgb = SLATE_BLUE
        card.line.color.rgb = EMERALD

        tf_c = card.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = Inches(0.3)
        tf_c.margin_top = Inches(0.3)

        p_head = tf_c.paragraphs[0]
        p_head.text = p_title
        p_head.font.size = Pt(18)
        p_head.font.bold = True
        p_head.font.color.rgb = EMERALD

        p_desc = tf_c.add_paragraph()
        p_desc.text = p_body
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = RGBColor(226, 232, 240)
        p_desc.space_before = Pt(14)

    # Output PowerPoint deck path
    output_pptx = Path("retsim_financial_planner_presentation.pptx").resolve()
    prs.save(output_pptx)
    print(f"Presentation saved successfully to {output_pptx}")

if __name__ == "__main__":
    create_deck()
